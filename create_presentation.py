#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
import os

def read_thesis_content(docx_path):
    """Extract content from thesis document"""
    doc = Document(docx_path)
    content = {}
    current_section = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            if para.style.name.startswith('Heading'):
                current_section = text
                content[current_section] = []
            elif current_section:
                content[current_section].append(text)

    return content

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_items, layout_type=1):
    """Add content slide with bullet points"""
    slide_layout = prs.slide_layouts[layout_type]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Add content
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(14)

    return slide

def add_image_slide(prs, title, image_path):
    """Add slide with image"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title

    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(32)

    # Add image if exists
    if os.path.exists(image_path):
        left = Inches(1.5)
        top = Inches(1.5)
        slide.shapes.add_picture(image_path, left, top, height=Inches(4.5))

    return slide

def create_presentation():
    """Main function to create presentation"""

    # Paths
    thesis_path = r"C:\Users\Iurii\Desktop\magister\ФКНТ_2025_121_магістр_Хоменко Ю.Ю..docx"
    output_path = r"C:\Users\Iurii\Desktop\magister\Презентація_Хоменко.pptx"

    # Diagram paths
    diagrams_base = r"C:\Users\Iurii\Desktop\magister\diagrams"
    use_case_diagram = os.path.join(diagrams_base, "use-case-diagram.drawio.png")  # May not exist
    class_diagram = os.path.join(diagrams_base, "class_diagram", "class-diagram-v2.drawio.png")
    package_diagram = os.path.join(diagrams_base, "package_diagram", "package-diagram.drawio.png")
    db_diagram = os.path.join(diagrams_base, "database_conceptual.drawio.png")

    # Check which diagrams exist
    diagrams_exist = {
        'use_case': os.path.exists(use_case_diagram),
        'class': os.path.exists(class_diagram),
        'package': os.path.exists(package_diagram),
        'database': os.path.exists(db_diagram)
    }

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Багатоплатформенний застосунок магазину музичних інструментів",
        "Виконав: студент групи ПІ-521м Хоменко Ю.Ю.\nКерівник: ст.в. Яцко К.С."
    )

    # Slide 2: Мета, об'єкт, предмет, завдання
    add_content_slide(
        prs,
        "Мета та завдання проєкту",
        [
            "Мета: розробка багатоплатформенного програмного засобу для онлайн-магазину",
            "Об'єкт: забезпечення можливості організації багатоплатформенної системи електронної комерції",
            "Предмет: структура та реалізація багатоплатформенного застосунку для онлайн-магазину",
            "",
            "Завдання:",
            "• Провести аналіз термінологічного апарату дослідження",
            "• Визначити специфіку та вимоги багатоплатформенних систем онлайн-комерції",
            "• Дослідити методики розробки ПЗ та обрати оптимальну архітектуру",
            "• Спроектувати та реалізувати прототип системи",
            "• Провести тестування та оцінку якості розробленого програмного засобу"
        ]
    )

    # Slide 3: Актуальність
    add_content_slide(
        prs,
        "Актуальність дослідження",
        [
            "Зростання електронної комерції:",
            "• ≈63.3% глобального веб-трафіку йде з мобільних пристроїв",
            "• ≈57% світових продажів у 2024 році через мобільні пристрої",
            "• 75% відвідувань сайтів ритейлерів з смартфонів",
            "",
            "Потреба в багатоплатформенних рішеннях:",
            "• Охоплення широкої аудиторії",
            "• Зручний доступ користувачів незалежно від пристрою",
            "• Конкурентоспроможність на ринку"
        ]
    )

    # Slide 4: Аналіз конкурентів
    add_content_slide(
        prs,
        "Аналіз конкурентів",
        [
            "Проаналізовані платформи:",
            "• Sweetwater - найбільший онлайн-магазин музичних інструментів у США",
            "• Guitar Center - мережа роздрібних магазинів з онлайн-платформою",
            "• Thomann - європейський лідер у сфері продажу музичного обладнання",
            "",
            "Ключові особливості конкурентів:",
            "• Широкий асортимент продукції",
            "• Системи рекомендацій та відгуків",
            "• Інтеграція з платіжними системами",
            "• Мобільні застосунки та адаптивний веб-дизайн"
        ]
    )

    # Slide 5: Архітектура системи
    add_content_slide(
        prs,
        "Архітектура системи",
        [
            "Обрано клієнт-серверну архітектуру з мікросервісним підходом",
            "",
            "Компоненти системи:",
            "• Web-застосунок (Next.js + React)",
            "• Mobile-застосунок (React Native)",
            "• Backend API (NestJS + Node.js)",
            "• База даних (PostgreSQL)",
            "• Спільні модулі (Common Package)",
            "",
            "Переваги архітектури:",
            "• Масштабованість та гнучкість",
            "• Модульність та повторне використання коду",
            "• Незалежний розвиток компонентів"
        ]
    )

    # Slide 6: Технології
    add_content_slide(
        prs,
        "Інструменти та технології",
        [
            "Frontend:",
            "• React, Next.js - веб-застосунок",
            "• React Native - мобільний застосунок",
            "• TypeScript - типізація коду",
            "",
            "Backend:",
            "• NestJS - фреймворк для серверної частини",
            "• PostgreSQL - реляційна база даних",
            "• TypeORM - ORM для роботи з БД",
            "",
            "DevOps та інструменти:",
            "• Docker - контейнеризація",
            "• PNPM - пакетний менеджер",
            "• Git - система контролю версій"
        ]
    )

    # Slide 7: Use Case Diagram
    if os.path.exists(use_case_diagram):
        add_image_slide(prs, "Діаграма варіантів використання", use_case_diagram)
    else:
        add_content_slide(
            prs,
            "Діаграма варіантів використання",
            [
                "Основні актори:",
                "• Користувач (покупець)",
                "• Адміністратор",
                "",
                "Функціональність для користувача:",
                "• Перегляд каталогу товарів",
                "• Пошук та фільтрація",
                "• Управління кошиком",
                "• Оформлення замовлення",
                "• Перегляд історії замовлень",
                "",
                "Функціональність для адміністратора:",
                "• Управління товарами",
                "• Обробка замовлень",
                "• Управління користувачами"
            ]
        )

    # Slide 8: Class Diagram
    if os.path.exists(class_diagram):
        add_image_slide(prs, "Діаграма класів", class_diagram)
    else:
        add_content_slide(
            prs,
            "Діаграма класів",
            [
                "Основні класи системи:",
                "• User - користувач системи",
                "• Product - товар",
                "• Order - замовлення",
                "• Cart - кошик",
                "• Review - відгук",
                "• Category - категорія товарів",
                "",
                "Зв'язки між класами:",
                "• Один користувач може мати багато замовлень",
                "• Замовлення містить багато товарів",
                "• Товар належить до категорії"
            ]
        )

    # Slide 9: Package Diagram
    if os.path.exists(package_diagram):
        add_image_slide(prs, "Діаграма пакетів", package_diagram)
    else:
        add_content_slide(
            prs,
            "Діаграма пакетів",
            [
                "Структура проєкту (monorepo):",
                "",
                "apps/",
                "• web - веб-застосунок (Next.js)",
                "• mobile - мобільний застосунок (React Native)",
                "",
                "services/",
                "• api - серверна частина (NestJS)",
                "",
                "packages/",
                "• common - спільні типи, утиліти, валідація",
                "",
                "docs/",
                "• diagrams - UML діаграми"
            ]
        )

    # Slide 10: Database Diagram
    if os.path.exists(db_diagram):
        add_image_slide(prs, "Концептуальна модель бази даних", db_diagram)
    else:
        add_content_slide(
            prs,
            "Структура бази даних",
            [
                "Основні таблиці:",
                "• users - користувачі",
                "• products - товари",
                "• categories - категорії",
                "• orders - замовлення",
                "• order_items - елементи замовлень",
                "• cart_items - товари в кошику",
                "• reviews - відгуки",
                "• wishlist - список бажань",
                "",
                "Додаткові таблиці:",
                "• product_images - зображення товарів",
                "• addresses - адреси доставки"
            ]
        )

    # Slide 11-13: Прототип системи
    add_content_slide(
        prs,
        "Функціональність системи - Каталог товарів",
        [
            "Веб та мобільний застосунок включають:",
            "",
            "Каталог товарів:",
            "• Відображення товарів з фото, назвою, ціною",
            "• Багаторівнева система категорій",
            "• Фільтрація за категоріями, ціною, брендом",
            "• Пошук за ключовими словами",
            "• Сортування (популярність, ціна, новизна)",
            "",
            "Деталі товару:",
            "• Повний опис та характеристики",
            "• Галерея зображень",
            "• Відгуки та рейтинги користувачів",
            "• Рекомендовані товари"
        ]
    )

    add_content_slide(
        prs,
        "Функціональність системи - Кошик та замовлення",
        [
            "Кошик покупок:",
            "• Додавання/видалення товарів",
            "• Зміна кількості",
            "• Автоматичний розрахунок суми",
            "• Збереження кошика для авторизованих користувачів",
            "",
            "Оформлення замовлення:",
            "• Вибір адреси доставки",
            "• Вибір способу оплати",
            "• Підтвердження замовлення",
            "• Відстеження статусу замовлення",
            "",
            "Особистий кабінет:",
            "• Профіль користувача",
            "• Історія замовлень",
            "• Список бажань"
        ]
    )

    add_content_slide(
        prs,
        "Функціональність системи - Адміністрування",
        [
            "Панель адміністратора:",
            "",
            "Управління товарами:",
            "• Додавання/редагування/видалення товарів",
            "• Управління категоріями",
            "• Завантаження зображень",
            "• Контроль залишків на складі",
            "",
            "Управління замовленнями:",
            "• Перегляд всіх замовлень",
            "• Зміна статусу замовлення",
            "• Обробка оплат",
            "",
            "Аналітика:",
            "• Статистика продажів",
            "• Популярні товари",
            "• Активність користувачів"
        ]
    )

    # Slide 14: Висновки
    add_content_slide(
        prs,
        "Висновки",
        [
            "Результати проєкту:",
            "",
            "• Проведено аналіз ринку конкурентів та сформовано вимоги до системи",
            "",
            "• Спроектовано клієнт-серверну архітектуру з мікросервісним підходом",
            "",
            "• Обрано сучасний стек технологій для багатоплатформенної розробки",
            "",
            "• Реалізовано прототип системи з функціональністю:",
            "  - Каталог товарів з пошуком та фільтрацією",
            "  - Кошик та оформлення замовлень",
            "  - Особистий кабінет користувача",
            "  - Панель адміністрування",
            "",
            "• Забезпечено можливість масштабування та подальшого розвитку системи"
        ]
    )

    # Save presentation
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    try:
        output = create_presentation()
        print("SUCCESS")
    except Exception as e:
        print("ERROR")
        import traceback
        traceback.print_exc()
