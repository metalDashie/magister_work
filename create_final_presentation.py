#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)

    return slide

def add_content_slide(prs, title, content_items):
    """Add content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for item in content_items:
        if isinstance(item, dict):
            p = tf.add_paragraph()
            p.text = item['text']
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('size', 16))
            if item.get('bold'):
                p.font.bold = True
        else:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(16)

    return slide

def add_image_slide(prs, title, image_path, caption=None):
    """Add slide with image and optional caption"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(31, 73, 125)
    p.alignment = PP_ALIGN.CENTER

    # Add image
    if os.path.exists(image_path):
        left = Inches(1)
        top = Inches(1.2)
        max_width = Inches(8)
        max_height = Inches(5.5)

        pic = slide.shapes.add_picture(image_path, left, top, width=max_width)

        # Center the image
        left = (prs.slide_width - pic.width) / 2
        pic.left = int(left)

    # Add caption if provided
    if caption:
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(9)
        height = Inches(0.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = caption
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_two_images_slide(prs, title, image1_path, image2_path, caption1=None, caption2=None):
    """Add slide with two images side by side"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(31, 73, 125)
    p.alignment = PP_ALIGN.CENTER

    # Add first image (left)
    if os.path.exists(image1_path):
        left = Inches(0.5)
        top = Inches(1.3)
        max_height = Inches(5)

        pic1 = slide.shapes.add_picture(image1_path, left, top, height=max_height)

    # Add second image (right)
    if os.path.exists(image2_path):
        left = Inches(5.5)
        top = Inches(1.3)
        max_height = Inches(5)

        pic2 = slide.shapes.add_picture(image2_path, left, top, height=max_height)

    return slide

def create_presentation():
    """Create final presentation with screenshots"""

    output_path = r"C:\Users\Iurii\Desktop\magister\Презентація_Хоменко_Фінальна.pptx"

    # Image paths
    extracted_base = r"C:\Users\Iurii\Desktop\magister\extracted_images"
    diagrams_base = r"C:\Users\Iurii\Desktop\magister\diagrams"

    tech_stack_img = os.path.join(extracted_base, "image12.png")
    web_mobile_screens = os.path.join(extracted_base, "image15.png")
    web_catalog = os.path.join(extracted_base, "image25.png")
    mobile_catalog = os.path.join(extracted_base, "image43.png")
    mobile_screen2 = os.path.join(extracted_base, "image44.png")
    mobile_screen3 = os.path.join(extracted_base, "image46.png")

    class_diagram = os.path.join(diagrams_base, "class_diagram", "class-diagram-v2.drawio.png")
    package_diagram = os.path.join(diagrams_base, "package_diagram", "package-diagram.drawio.png")
    db_diagram = os.path.join(diagrams_base, "database_conceptual.drawio.png")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Багатоплатформенний застосунок\nмагазину музичних інструментів",
        "Виконав: студент групи ПІ-521м Хоменко Ю.Ю.\nКерівник: ст.в. Яцко К.С."
    )

    # Slide 2: Мета та завдання
    add_content_slide(
        prs,
        "Мета та завдання проєкту",
        [
            {'text': 'Мета проєкту:', 'bold': True, 'size': 18},
            {'text': 'Розробка багатоплатформенного програмного засобу для онлайн-магазину з дослідженням методик і варіантів розробки ПЗ', 'level': 1, 'size': 15},
            '',
            {'text': 'Об\'єкт дослідження:', 'bold': True, 'size': 16},
            {'text': 'Забезпечення можливості організації багатоплатформенної системи електронної комерції', 'level': 1, 'size': 14},
            '',
            {'text': 'Предмет дослідження:', 'bold': True, 'size': 16},
            {'text': 'Структура та реалізація багатоплатформенного застосунку для онлайн-магазину', 'level': 1, 'size': 14},
        ]
    )

    # Slide 3: Завдання
    add_content_slide(
        prs,
        "Завдання дослідження",
        [
            "Проаналізувати термінологічний апарат дослідження",
            "Визначити специфіку, вимоги та особливості багатоплатформенних систем",
            "Дослідити методики розробки ПЗ та обрати оптимальну",
            "Провести аналіз ринку конкурентів",
            "Спроектувати архітектуру програмної системи",
            "Обрати технологічний стек для розробки",
            "Реалізувати прототип системи з основною функціональністю",
            "Провести тестування та оцінку якості ПЗ"
        ]
    )

    # Slide 4: Актуальність
    add_content_slide(
        prs,
        "Актуальність дослідження",
        [
            {'text': 'Зростання електронної комерції:', 'bold': True, 'size': 18},
            {'text': '≈63.3% глобального веб-трафіку йде з мобільних пристроїв', 'level': 1},
            {'text': '≈57% світових продажів у 2024 році через мобільні пристрої', 'level': 1},
            {'text': '75% відвідувань сайтів ритейлерів зі смартфонів', 'level': 1},
            '',
            {'text': 'Необхідність багатоплатформенних рішень:', 'bold': True, 'size': 18},
            {'text': 'Охоплення широкої аудиторії (Web + iOS + Android)', 'level': 1},
            {'text': 'Зручний доступ незалежно від пристрою', 'level': 1},
            {'text': 'Підвищення конкурентоспроможності', 'level': 1},
        ]
    )

    # Slide 5: Аналіз конкурентів
    add_content_slide(
        prs,
        "Аналіз конкурентів",
        [
            {'text': 'Sweetwater (США)', 'bold': True, 'size': 18},
            {'text': 'Найбільший онлайн-магазин музичних інструментів', 'level': 1, 'size': 14},
            {'text': 'Детальні описи, відео-демонстрації', 'level': 1, 'size': 14},
            '',
            {'text': 'Thomann (Німеччина)', 'bold': True, 'size': 18},
            {'text': 'Європейський лідер, понад 1 млн продуктів', 'level': 1, 'size': 14},
            {'text': 'Багатомовність, розширена фільтрація', 'level': 1, 'size': 14},
            '',
            {'text': 'Guitar Center (США)', 'bold': True, 'size': 18},
            {'text': 'Мережа магазинів + онлайн-платформа', 'level': 1, 'size': 14},
            {'text': 'Програма лояльності, обмін інструментів', 'level': 1, 'size': 14},
        ]
    )

    # Slide 6: Методологія Agile/Scrum
    add_content_slide(
        prs,
        "Методологія розробки: Agile/Scrum",
        [
            {'text': 'Agile - гнучка методологія розробки:', 'bold': True, 'size': 18},
            {'text': 'Короткі ітерації (спринти)', 'level': 1},
            {'text': 'Тісна взаємодія з замовником', 'level': 1},
            {'text': 'Швидке реагування на зміни', 'level': 1},
            '',
            {'text': 'Scrum - структура роботи:', 'bold': True, 'size': 18},
            {'text': 'Спринти тривалістю 2-4 тижні', 'level': 1},
            {'text': 'Ефективне планування та коригування', 'level': 1},
            {'text': 'Регулярні retrospective', 'level': 1},
        ]
    )

    # Slide 7: Архітектура
    add_content_slide(
        prs,
        "Архітектура системи",
        [
            {'text': 'Клієнт-серверна архітектура з елементами мікросервісів', 'bold': True, 'size': 18},
            '',
            {'text': 'Клієнтська частина:', 'bold': True},
            {'text': 'Web-застосунок (Next.js + React)', 'level': 1},
            {'text': 'Mobile-застосунок (React Native)', 'level': 1},
            '',
            {'text': 'Серверна частина:', 'bold': True},
            {'text': 'RESTful API (NestJS + Node.js)', 'level': 1},
            {'text': 'PostgreSQL з TypeORM', 'level': 1},
            '',
            {'text': 'Спільні компоненти:', 'bold': True},
            {'text': 'Common Package - типи, валідація, утиліти', 'level': 1},
        ]
    )

    # Slide 8: Technology Stack (з картинкою)
    if os.path.exists(tech_stack_img):
        add_image_slide(
            prs,
            "Технологічний стек системи",
            tech_stack_img,
            "Frontend (Web/Mobile), Backend (API), Payments & Integrations, DevTools & Infrastructure"
        )
    else:
        add_content_slide(
            prs,
            "Технології Frontend",
            [
                {'text': 'React + TypeScript', 'bold': True},
                {'text': 'Компонентний підхід, типобезпека', 'level': 1, 'size': 14},
                '',
                {'text': 'Next.js (Web)', 'bold': True},
                {'text': 'SSR/SSG, оптимізація SEO', 'level': 1, 'size': 14},
                '',
                {'text': 'React Native (Mobile)', 'bold': True},
                {'text': 'Єдина кодова база для iOS/Android', 'level': 1, 'size': 14},
            ]
        )

    # Slide 9: Package Diagram
    if os.path.exists(package_diagram):
        add_image_slide(
            prs,
            "Діаграма пакетів",
            package_diagram,
            "Monorepo структура проєкту"
        )

    # Slide 10: Class Diagram
    if os.path.exists(class_diagram):
        add_image_slide(
            prs,
            "Діаграма класів",
            class_diagram,
            "Основні сутності системи"
        )

    # Slide 11: Database Diagram
    if os.path.exists(db_diagram):
        add_image_slide(
            prs,
            "Концептуальна модель бази даних",
            db_diagram,
            "Структура таблиць БД"
        )

    # Slide 12: Прототип - Web і Mobile огляд
    if os.path.exists(web_mobile_screens):
        add_image_slide(
            prs,
            "Прототип системи: Web та Mobile",
            web_mobile_screens,
            "Веб-застосунок та мобільний застосунок"
        )

    # Slide 13: Веб-застосунок - Каталог
    if os.path.exists(web_catalog):
        add_image_slide(
            prs,
            "Веб-застосунок: Каталог товарів",
            web_catalog,
            "Перегляд товарів з фільтрацією, пошуком та сортуванням"
        )

    # Slide 14: Мобільний застосунок - Каталог
    if os.path.exists(mobile_catalog):
        add_image_slide(
            prs,
            "Мобільний застосунок: Каталог",
            mobile_catalog,
            "Мобільна версія каталогу з адаптивним інтерфейсом"
        )

    # Slide 15: Мобільні екрани 2 і 3
    if os.path.exists(mobile_screen2) and os.path.exists(mobile_screen3):
        add_two_images_slide(
            prs,
            "Мобільний застосунок: Додаткові екрани",
            mobile_screen2,
            mobile_screen3
        )

    # Slide 16: Функціональність
    add_content_slide(
        prs,
        "Основна функціональність системи",
        [
            {'text': 'Для користувачів:', 'bold': True, 'size': 18},
            {'text': 'Каталог товарів з пошуком та фільтрацією', 'level': 1},
            {'text': 'Детальна інформація про товар', 'level': 1},
            {'text': 'Кошик та оформлення замовлень', 'level': 1},
            {'text': 'Відгуки та рейтинги товарів', 'level': 1},
            {'text': 'Особистий кабінет та історія замовлень', 'level': 1},
            '',
            {'text': 'Для адміністраторів:', 'bold': True, 'size': 18},
            {'text': 'Управління товарами та категоріями', 'level': 1},
            {'text': 'Обробка замовлень', 'level': 1},
            {'text': 'Модерація відгуків', 'level': 1},
        ]
    )

    # Slide 17: Безпека
    add_content_slide(
        prs,
        "Безпека системи",
        [
            {'text': 'Автентифікація та авторизація:', 'bold': True},
            {'text': 'JWT токени з refresh механізмом', 'level': 1, 'size': 14},
            {'text': 'Role-based access control (User, Admin)', 'level': 1, 'size': 14},
            '',
            {'text': 'Захист даних:', 'bold': True},
            {'text': 'Bcrypt для хешування паролів', 'level': 1, 'size': 14},
            {'text': 'HTTPS для всіх з\'єднань', 'level': 1, 'size': 14},
            {'text': 'Валідація даних (class-validator)', 'level': 1, 'size': 14},
            '',
            {'text': 'Захист від атак:', 'bold': True},
            {'text': 'CORS налаштування', 'level': 1, 'size': 14},
            {'text': 'Rate limiting для API', 'level': 1, 'size': 14},
        ]
    )

    # Slide 18: Висновки
    add_content_slide(
        prs,
        "Висновки",
        [
            {'text': 'Досягнуті результати:', 'bold': True, 'size': 18},
            '',
            "Проаналізовано методи дослідження та принципи UI/UX",
            "Досліджено ринок та сформовано вимоги до системи",
            "Обрано методологію Agile/Scrum",
            "Спроектовано масштабовану архітектуру",
            "Обрано сучасний технологічний стек",
            "Реалізовано прототип з веб та мобільними застосунками",
            "Впроваджено систему безпеки та автентифікації",
            "Забезпечено можливість подальшого розвитку"
        ]
    )

    # Save presentation
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    try:
        output = create_presentation()
        print("SUCCESS")
    except Exception as e:
        print("ERROR")
        import traceback
        traceback.print_exc()
