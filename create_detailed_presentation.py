#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from docx import Document
import os

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title

    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)

    return slide

def add_content_slide(prs, title, content_items):
    """Add content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for item in content_items:
        if isinstance(item, dict):
            p = tf.add_paragraph()
            p.text = item['text']
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('size', 16))
            if item.get('bold'):
                p.font.bold = True
        else:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(16)

    return slide

def add_image_slide(prs, title, image_path, caption=None):
    """Add slide with image and optional caption"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(31, 73, 125)
    p.alignment = PP_ALIGN.CENTER

    # Add image
    if os.path.exists(image_path):
        left = Inches(1)
        top = Inches(1.2)
        max_width = Inches(8)
        max_height = Inches(5.5)

        pic = slide.shapes.add_picture(image_path, left, top, width=max_width)

        # Center the image
        left = (prs.slide_width - pic.width) / 2
        pic.left = int(left)

    # Add caption if provided
    if caption:
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(9)
        height = Inches(0.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = caption
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation():
    """Create detailed presentation"""

    output_path = r"C:\Users\Iurii\Desktop\magister\Презентація_Хоменко_Детальна.pptx"

    # Diagram paths
    diagrams_base = r"C:\Users\Iurii\Desktop\magister\diagrams"
    class_diagram = os.path.join(diagrams_base, "class_diagram", "class-diagram-v2.drawio.png")
    package_diagram = os.path.join(diagrams_base, "package_diagram", "package-diagram.drawio.png")
    db_diagram = os.path.join(diagrams_base, "database_conceptual.drawio.png")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Багатоплатформенний застосунок\nмагазину музичних інструментів",
        "Виконав: студент групи ПІ-521м Хоменко Ю.Ю.\nКерівник: ст.в. Яцко К.С."
    )

    # Slide 2: Мета та завдання
    add_content_slide(
        prs,
        "Мета та завдання проєкту",
        [
            {'text': 'Мета проєкту:', 'bold': True, 'size': 18},
            {'text': 'Розробка багатоплатформенного програмного засобу для онлайн-магазину з дослідженням методик і варіантів розробки ПЗ', 'level': 1, 'size': 15},
            '',
            {'text': 'Об\'єкт дослідження:', 'bold': True, 'size': 16},
            {'text': 'Забезпечення можливості організації багатоплатформенної системи електронної комерції', 'level': 1, 'size': 14},
            '',
            {'text': 'Предмет дослідження:', 'bold': True, 'size': 16},
            {'text': 'Структура та реалізація багатоплатформенного застосунку для онлайн-магазину', 'level': 1, 'size': 14},
        ]
    )

    # Slide 3: Завдання дослідження
    add_content_slide(
        prs,
        "Завдання дослідження",
        [
            "Проаналізувати термінологічний апарат дослідження",
            "Визначити специфіку, вимоги та характерні особливості багатоплатформенних систем онлайн-комерції",
            "Дослідити методики розробки ПЗ та обрати оптимальну",
            "Провести аналіз ринку конкурентів та сформувати вимоги",
            "Спроектувати архітектуру програмної системи",
            "Обрати технологічний стек для розробки",
            "Реалізувати прототип системи з основною функціональністю",
            "Провести тестування та оцінку якості ПЗ"
        ]
    )

    # Slide 4: Актуальність
    add_content_slide(
        prs,
        "Актуальність дослідження",
        [
            {'text': 'Зростання електронної комерції:', 'bold': True, 'size': 18},
            {'text': '≈63.3% глобального веб-трафіку йде з мобільних пристроїв (StatCounter)', 'level': 1},
            {'text': '≈57% світових продажів у 2024 році через мобільні пристрої', 'level': 1},
            {'text': '75% відвідувань сайтів ритейлерів здійснюються зі смартфонів', 'level': 1},
            '',
            {'text': 'Необхідність багатоплатформенних рішень:', 'bold': True, 'size': 18},
            {'text': 'Охоплення широкої аудиторії (Web + iOS + Android)', 'level': 1},
            {'text': 'Забезпечення зручного доступу незалежно від пристрою', 'level': 1},
            {'text': 'Підвищення конкурентоспроможності на ринку', 'level': 1},
        ]
    )

    # Slide 5: Методи дослідження
    add_content_slide(
        prs,
        "Методи дослідження аудиторії",
        [
            {'text': 'Формування портрету ідеального покупця (User Persona)', 'level': 0},
            {'text': 'Визначення потреб, очікувань та поведінкових патернів', 'level': 1, 'size': 14},
            '',
            {'text': 'Інтерв\'ю та фокус-групи', 'level': 0},
            {'text': 'Якісні дані про ставлення до інтерфейсу та очікування', 'level': 1, 'size': 14},
            '',
            {'text': 'Аналіз поведінкових даних (Behavioral Analytics)', 'level': 0},
            {'text': 'Вивчення натискань, часу на сторінці, конверсій', 'level': 1, 'size': 14},
            '',
            {'text': 'A/B-тестування', 'level': 0},
            {'text': 'Порівняння версій інтерфейсу для визначення ефективності', 'level': 1, 'size': 14},
        ]
    )

    # Slide 6: UI/UX принципи
    add_content_slide(
        prs,
        "Основні закони створення інтерфейсів",
        [
            {'text': 'Закон Якоба (Jakob\'s Law):', 'bold': True},
            {'text': 'Користувачі очікують знайомих паттернів інтерфейсу', 'level': 1, 'size': 14},
            '',
            {'text': 'Закон Фітта (Fitts\'s Law):', 'bold': True},
            {'text': 'Час досягнення елемента залежить від відстані та розміру', 'level': 1, 'size': 14},
            '',
            {'text': 'Закон Хікса (Hick\'s Law):', 'bold': True},
            {'text': 'Час прийняття рішення зростає з кількістю варіантів', 'level': 1, 'size': 14},
            '',
            {'text': 'Закон Міллера (Miller\'s Law):', 'bold': True},
            {'text': 'Людина може утримувати 7±2 елементи в короткочасній пам\'яті', 'level': 1, 'size': 14},
        ]
    )

    # Slide 7: Аналіз конкурентів
    add_content_slide(
        prs,
        "Аналіз конкурентів",
        [
            {'text': 'Sweetwater (США)', 'bold': True, 'size': 18},
            {'text': 'Найбільший онлайн-магазин музичних інструментів', 'level': 1, 'size': 14},
            {'text': 'Детальні описи продуктів, відео-демонстрації', 'level': 1, 'size': 14},
            '',
            {'text': 'Thomann (Німеччина)', 'bold': True, 'size': 18},
            {'text': 'Європейський лідер, понад 1 млн продуктів', 'level': 1, 'size': 14},
            {'text': 'Багатомовність, розширена система фільтрації', 'level': 1, 'size': 14},
            '',
            {'text': 'Guitar Center (США)', 'bold': True, 'size': 18},
            {'text': 'Мережа роздрібних магазинів + онлайн-платформа', 'level': 1, 'size': 14},
            {'text': 'Програма лояльності, можливість обміну інструментів', 'level': 1, 'size': 14},
        ]
    )

    # Slide 8: Вибір методології
    add_content_slide(
        prs,
        "Методологія розробки: Agile/Scrum",
        [
            {'text': 'Agile - гнучка методологія розробки:', 'bold': True, 'size': 18},
            {'text': 'Короткі ітерації (спринти)', 'level': 1},
            {'text': 'Тісна взаємодія з замовником', 'level': 1},
            {'text': 'Швидке реагування на зміни', 'level': 1},
            '',
            {'text': 'Scrum - структура роботи:', 'bold': True, 'size': 18},
            {'text': 'Спринти тривалістю 2-4 тижні', 'level': 1},
            {'text': 'Ефективне планування та коригування задач', 'level': 1},
            {'text': 'Регулярні retrospective для покращення процесів', 'level': 1},
        ]
    )

    # Slide 9: Архітектура
    add_content_slide(
        prs,
        "Архітектура системи",
        [
            {'text': 'Клієнт-серверна архітектура з елементами мікросервісів', 'bold': True, 'size': 18},
            '',
            {'text': 'Клієнтська частина:', 'bold': True},
            {'text': 'Web-застосунок (Next.js + React)', 'level': 1},
            {'text': 'Mobile-застосунок (React Native)', 'level': 1},
            '',
            {'text': 'Серверна частина:', 'bold': True},
            {'text': 'RESTful API (NestJS + Node.js)', 'level': 1},
            {'text': 'PostgreSQL з TypeORM', 'level': 1},
            '',
            {'text': 'Спільні компоненти:', 'bold': True},
            {'text': 'Common Package - типи, валідація, утиліти', 'level': 1},
        ]
    )

    # Slide 10: Технології Frontend
    add_content_slide(
        prs,
        "Технології: Frontend",
        [
            {'text': 'React + TypeScript', 'bold': True},
            {'text': 'Компонентний підхід, типобезпека', 'level': 1, 'size': 14},
            '',
            {'text': 'Next.js (Web)', 'bold': True},
            {'text': 'SSR/SSG, оптимізація SEO, App Router', 'level': 1, 'size': 14},
            '',
            {'text': 'React Native (Mobile)', 'bold': True},
            {'text': 'Єдина кодова база для iOS/Android', 'level': 1, 'size': 14},
            '',
            {'text': 'Tailwind CSS', 'bold': True},
            {'text': 'Швидка розробка адаптивного дизайну', 'level': 1, 'size': 14},
            '',
            {'text': 'React Query', 'bold': True},
            {'text': 'Управління серверним станом, кешування', 'level': 1, 'size': 14},
        ]
    )

    # Slide 11: Технології Backend
    add_content_slide(
        prs,
        "Технології: Backend",
        [
            {'text': 'NestJS', 'bold': True},
            {'text': 'Модульна архітектура, Dependency Injection', 'level': 1, 'size': 14},
            '',
            {'text': 'PostgreSQL', 'bold': True},
            {'text': 'Надійна реляційна СУБД з підтримкою ACID', 'level': 1, 'size': 14},
            '',
            {'text': 'TypeORM', 'bold': True},
            {'text': 'ORM для типобезпечної роботи з БД', 'level': 1, 'size': 14},
            '',
            {'text': 'Passport.js + JWT', 'bold': True},
            {'text': 'Автентифікація та авторизація користувачів', 'level': 1, 'size': 14},
            '',
            {'text': 'Docker', 'bold': True},
            {'text': 'Контейнеризація для консистентного середовища', 'level': 1, 'size': 14},
        ]
    )

    # Slide 12: Package Diagram
    if os.path.exists(package_diagram):
        add_image_slide(
            prs,
            "Діаграма пакетів (Структура проєкту)",
            package_diagram,
            "Monorepo структура з розділенням на клієнтські застосунки, серверну частину та спільні модулі"
        )
    else:
        add_content_slide(
            prs,
            "Структура проєкту (Monorepo)",
            [
                {'text': 'apps/', 'bold': True, 'size': 18},
                {'text': 'web - веб-застосунок (Next.js)', 'level': 1},
                {'text': 'mobile - мобільний застосунок (React Native)', 'level': 1},
                '',
                {'text': 'services/', 'bold': True, 'size': 18},
                {'text': 'api - серверна частина (NestJS)', 'level': 1},
                '',
                {'text': 'packages/', 'bold': True, 'size': 18},
                {'text': 'common - типи, валідація, утиліти', 'level': 1},
                '',
                {'text': 'docs/', 'bold': True, 'size': 18},
                {'text': 'diagrams - UML діаграми системи', 'level': 1},
            ]
        )

    # Slide 13: Class Diagram
    if os.path.exists(class_diagram):
        add_image_slide(
            prs,
            "Діаграма класів",
            class_diagram,
            "Основні сутності системи та їх взаємозв'язки"
        )

    # Slide 14: Database Diagram
    if os.path.exists(db_diagram):
        add_image_slide(
            prs,
            "Концептуальна модель бази даних",
            db_diagram,
            "Структура таблиць БД з відношеннями між сутностями"
        )

    # Slide 15: Функціональність - Каталог
    add_content_slide(
        prs,
        "Функціональність: Каталог товарів",
        [
            {'text': 'Перегляд товарів', 'bold': True},
            {'text': 'Відображення з фото, назвою, ціною, рейтингом', 'level': 1, 'size': 14},
            '',
            {'text': 'Система категорій', 'bold': True},
            {'text': 'Багаторівнева ієрархія (Гітари → Електрогітари → Stratocaster)', 'level': 1, 'size': 14},
            '',
            {'text': 'Фільтрація та пошук', 'bold': True},
            {'text': 'За категоріями, ціною, брендом, характеристиками', 'level': 1, 'size': 14},
            {'text': 'Full-text пошук за ключовими словами', 'level': 1, 'size': 14},
            '',
            {'text': 'Сортування', 'bold': True},
            {'text': 'За популярністю, ціною (зростання/спадання), новизною', 'level': 1, 'size': 14},
        ]
    )

    # Slide 16: Функціональність - Детальна сторінка
    add_content_slide(
        prs,
        "Функціональність: Деталі товару",
        [
            "Повний опис товару з технічними характеристиками",
            "Галерея зображень (до 10 фото на товар)",
            "Відгуки користувачів з рейтингом (1-5 зірок)",
            "Можливість додати до кошика або списку бажань",
            "Рекомендовані товари (на основі категорії)",
            "Інформація про наявність на складі",
            "Можливість залишити відгук після покупки",
        ]
    )

    # Slide 17: Функціональність - Кошик
    add_content_slide(
        prs,
        "Функціональність: Кошик та оформлення",
        [
            {'text': 'Кошик покупок:', 'bold': True, 'size': 18},
            {'text': 'Додавання/видалення товарів', 'level': 1},
            {'text': 'Зміна кількості з автоматичним перерахунком', 'level': 1},
            {'text': 'Збереження кошика для авторизованих користувачів', 'level': 1},
            '',
            {'text': 'Оформлення замовлення:', 'bold': True, 'size': 18},
            {'text': 'Введення/вибір адреси доставки', 'level': 1},
            {'text': 'Вибір способу оплати (карта, готівка при отриманні)', 'level': 1},
            {'text': 'Підтвердження замовлення з email-повідомленням', 'level': 1},
            {'text': 'Відстеження статусу замовлення', 'level': 1},
        ]
    )

    # Slide 18: Функціональність - Адмін
    add_content_slide(
        prs,
        "Функціональність: Панель адміністратора",
        [
            {'text': 'Управління товарами:', 'bold': True},
            {'text': 'CRUD операції для продуктів та категорій', 'level': 1, 'size': 14},
            {'text': 'Завантаження та управління зображеннями', 'level': 1, 'size': 14},
            '',
            {'text': 'Управління замовленнями:', 'bold': True},
            {'text': 'Перегляд та обробка замовлень', 'level': 1, 'size': 14},
            {'text': 'Зміна статусу (нове, в обробці, відправлено, доставлено)', 'level': 1, 'size': 14},
            '',
            {'text': 'Управління користувачами:', 'bold': True},
            {'text': 'Перегляд списку користувачів', 'level': 1, 'size': 14},
            {'text': 'Модерація відгуків', 'level': 1, 'size': 14},
        ]
    )

    # Slide 19: Безпека
    add_content_slide(
        prs,
        "Безпека системи",
        [
            {'text': 'Автентифікація та авторизація:', 'bold': True},
            {'text': 'JWT токени з refresh механізмом', 'level': 1, 'size': 14},
            {'text': 'Role-based access control (User, Admin)', 'level': 1, 'size': 14},
            '',
            {'text': 'Захист даних:', 'bold': True},
            {'text': 'Bcrypt для хешування паролів', 'level': 1, 'size': 14},
            {'text': 'HTTPS для всіх з\'єднань', 'level': 1, 'size': 14},
            {'text': 'Валідація даних на backend (class-validator)', 'level': 1, 'size': 14},
            '',
            {'text': 'Захист від атак:', 'bold': True},
            {'text': 'CORS налаштування', 'level': 1, 'size': 14},
            {'text': 'Rate limiting для API', 'level': 1, 'size': 14},
            {'text': 'SQL injection prevention (TypeORM)', 'level': 1, 'size': 14},
        ]
    )

    # Slide 20: Висновки
    add_content_slide(
        prs,
        "Висновки",
        [
            {'text': 'Досягнуті результати:', 'bold': True, 'size': 18},
            '',
            "Проаналізовано методи дослідження аудиторії та принципи UI/UX",
            "Досліджено ринок конкурентів та сформовано вимоги до системи",
            "Обрано методологію Agile/Scrum для гнучкої розробки",
            "Спроектовано масштабовану клієнт-серверну архітектуру",
            "Обрано сучасний стек: React, Next.js, React Native, NestJS",
            "Реалізовано повнофункціональний прототип з:",
            "  • Веб та мобільними застосунками",
            "  • Каталогом товарів з пошуком та фільтрацією",
            "  • Системою замовлень та оплати",
            "  • Панеллю адміністрування",
            "Забезпечено безпеку та можливість масштабування"
        ]
    )

    # Save presentation
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    try:
        output = create_presentation()
        print("SUCCESS")
    except Exception as e:
        print("ERROR")
        import traceback
        traceback.print_exc()
