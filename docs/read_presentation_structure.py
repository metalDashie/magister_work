#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from docx import Document
import sys

def read_presentation_structure(pptx_path):
    """Read and print structure of existing presentation"""
    print("=== СТРУКТУРА ІСНУЮЧОЇ ПРЕЗЕНТАЦІЇ ===\n")
    prs = Presentation(pptx_path)

    for i, slide in enumerate(prs.slides, 1):
        print(f"\n--- Слайд {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(f"  {shape.text[:100]}")

    return prs

def read_thesis_document(docx_path):
    """Read thesis document"""
    print("\n\n=== ЗМІСТ МАГІСТЕРСЬКОЇ РОБОТИ ===\n")
    doc = Document(docx_path)

    content = {}
    current_section = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Detect section headers
            if para.style.name.startswith('Heading'):
                current_section = text
                content[current_section] = []
                print(f"\n{text}")
            elif current_section:
                content[current_section].append(text)

    return content

if __name__ == "__main__":
    pptx_path = r"C:\Users\Iurii\Desktop\magister\дп_Презентація.pptx"
    docx_path = r"C:\Users\Iurii\Desktop\magister\ФКНТ_2025_121_магістр_Хоменко Ю.Ю..docx"

    try:
        prs = read_presentation_structure(pptx_path)
        print(f"\n\nВсього слайдів: {len(prs.slides)}")

        content = read_thesis_document(docx_path)
        print(f"\n\nВсього розділів: {len(content)}")

    except Exception as e:
        print(f"Помилка: {e}", file=sys.stderr)
        sys.exit(1)
