#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from docx import Document
import sys
import os

def analyze_presentation(pptx_path):
    """Analyze existing presentation structure"""
    result = []
    result.append("=== СТРУКТУРА ІСНУЮЧОЇ ПРЕЗЕНТАЦІЇ ===\n")

    try:
        prs = Presentation(pptx_path)
        result.append(f"Всього слайдів: {len(prs.slides)}\n")

        for i, slide in enumerate(prs.slides, 1):
            result.append(f"\n--- Слайд {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    result.append(f"  {shape.text}")

    except Exception as e:
        result.append(f"Помилка читання презентації: {e}")

    return result

def analyze_thesis(docx_path):
    """Analyze thesis document"""
    result = []
    result.append("\n\n=== ЗМІСТ МАГІСТЕРСЬКОЇ РОБОТИ ===\n")

    try:
        doc = Document(docx_path)
        result.append(f"Всього параграфів: {len(doc.paragraphs)}\n")

        # Get first 50 paragraphs to understand structure
        for i, para in enumerate(doc.paragraphs[:100], 1):
            text = para.text.strip()
            if text:
                style = para.style.name
                result.append(f"[{style}] {text[:150]}")

    except Exception as e:
        result.append(f"Помилка читання документа: {e}")

    return result

if __name__ == "__main__":
    pptx_path = r"C:\Users\Iurii\Desktop\magister\дп_Презентація.pptx"
    docx_path = r"C:\Users\Iurii\Desktop\magister\ФКНТ_2025_121_магістр_Хоменко Ю.Ю..docx"
    output_path = r"C:\magister_work\presentation_analysis.txt"

    all_results = []

    # Check if files exist
    if not os.path.exists(pptx_path):
        all_results.append(f"Файл не знайдено: {pptx_path}")
    else:
        all_results.extend(analyze_presentation(pptx_path))

    if not os.path.exists(docx_path):
        all_results.append(f"Файл не знайдено: {docx_path}")
    else:
        all_results.extend(analyze_thesis(docx_path))

    # Write to file
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(all_results))

    print(f"Results saved to: {output_path}")
    print(f"Total lines: {len(all_results)}")
