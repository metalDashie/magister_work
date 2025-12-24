#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)

    return slide

def add_content_slide(prs, title, content_items):
    """Add content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for item in content_items:
        if isinstance(item, dict):
            p = tf.add_paragraph()
            p.text = item['text']
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('size', 16))
            if item.get('bold'):
                p.font.bold = True
        else:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(16)

    return slide

def add_image_slide(prs, title, image_path, caption=None):
    """Add slide with image and optional caption"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(31, 73, 125)
    p.alignment = PP_ALIGN.CENTER

    # Add image
    if os.path.exists(image_path):
        left = Inches(1)
        top = Inches(1.2)
        max_width = Inches(8)

        pic = slide.shapes.add_picture(image_path, left, top, width=max_width)

        # Center the image
        left = (prs.slide_width - pic.width) / 2
        pic.left = int(left)

    # Add caption if provided
    if caption:
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(9)
        height = Inches(0.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = caption
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_two_images_slide(prs, title, image1_path, image2_path):
    """Add slide with two images side by side"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(31, 73, 125)
    p.alignment = PP_ALIGN.CENTER

    # Add first image (left)
    if os.path.exists(image1_path):
        left = Inches(0.5)
        top = Inches(1.3)
        max_height = Inches(5)

        pic1 = slide.shapes.add_picture(image1_path, left, top, height=max_height)

    # Add second image (right)
    if os.path.exists(image2_path):
        left = Inches(5.5)
        top = Inches(1.3)
        max_height = Inches(5)

        pic2 = slide.shapes.add_picture(image2_path, left, top, height=max_height)

    return slide

def create_presentation():
    """Create complete presentation with all sections"""

    output_path = r"C:\Users\Iurii\Desktop\magister\Презентація_Хоменко_Повна.pptx"

    # Image paths
    extracted_base = r"C:\Users\Iurii\Desktop\magister\extracted_images"
    diagrams_base = r"C:\Users\Iurii\Desktop\magister\diagrams"

    # Technology and screenshots
    tech_stack_img = os.path.join(extracted_base, "image12.png")
    web_mobile_screens = os.path.join(extracted_base, "image15.png")
    web_catalog = os.path.join(extracted_base, "image25.png")
    mobile_catalog = os.path.join(extracted_base, "image43.png")
    mobile_screen2 = os.path.join(extracted_base, "image44.png")
    mobile_screen3 = os.path.join(extracted_base, "image46.png")

    # Project management
    clickup_screenshot = os.path.join(extracted_base, "image13.png")
    scrum_diagram = os.path.join(extracted_base, "image7.png")

    # Admin panel
    user_management = os.path.join(extracted_base, "image36.png")
    analytics_dashboard = os.path.join(extracted_base, "image35.png")
    profile_screen = os.path.join(extracted_base, "image33.png")
    cart_screen = os.path.join(extracted_base, "image32.png")

    # Diagrams
    class_diagram = os.path.join(diagrams_base, "class_diagram", "class-diagram-v2.drawio.png")
    package_diagram = os.path.join(diagrams_base, "package_diagram", "package-diagram.drawio.png")
    db_diagram = os.path.join(diagrams_base, "database_conceptual.drawio.png")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==================== ВСТУП ====================

    # Slide 1: Title
    add_title_slide(
        prs,
        "Багатоплатформенний застосунок\nмагазину музичних інструментів",
        "Виконав: студент групи ПІ-521м Хоменко Ю.Ю.\nКерівник: ст.в. Яцко К.С."
    )

    # Slide 2: Мета та завдання
    add_content_slide(
        prs,
        "Мета та завдання проєкту",
        [
            {'text': 'Мета проєкту:', 'bold': True, 'size': 18},
            {'text': 'Розробка багатоплатформенного програмного засобу для онлайн-магазину з дослідженням методик і варіантів розробки ПЗ', 'level': 1, 'size': 15},
            '',
            {'text': 'Об\'єкт дослідження:', 'bold': True, 'size': 16},
            {'text': 'Забезпечення можливості організації багатоплатформенної системи електронної комерції', 'level': 1, 'size': 14},
            '',
            {'text': 'Предмет дослідження:', 'bold': True, 'size': 16},
            {'text': 'Структура та реалізація багатоплатформенного застосунку для онлайн-магазину', 'level': 1, 'size': 14},
        ]
    )

    # Slide 3: Завдання
    add_content_slide(
        prs,
        "Завдання дослідження",
        [
            "Проаналізувати термінологічний апарат дослідження",
            "Визначити специфіку та вимоги багатоплатформенних систем",
            "Дослідити методики розробки ПЗ та обрати оптимальну",
            "Провести аналіз ринку конкурентів",
            "Спроектувати архітектуру програмної системи",
            "Обрати технологічний стек",
            "Реалізувати прототип системи",
            "Провести тестування та оцінку якості"
        ]
    )

    # Slide 4: Актуальність
    add_content_slide(
        prs,
        "Актуальність дослідження",
        [
            {'text': 'Зростання електронної комерції:', 'bold': True, 'size': 18},
            {'text': '≈63.3% глобального веб-трафіку з мобільних пристроїв', 'level': 1},
            {'text': '≈57% світових продажів у 2024 році через мобільні', 'level': 1},
            {'text': '75% відвідувань сайтів ритейлерів зі смартфонів', 'level': 1},
            '',
            {'text': 'Необхідність багатоплатформенних рішень:', 'bold': True, 'size': 18},
            {'text': 'Охоплення широкої аудиторії', 'level': 1},
            {'text': 'Зручний доступ незалежно від пристрою', 'level': 1},
            {'text': 'Підвищення конкурентоспроможності', 'level': 1},
        ]
    )

    # Slide 5: Аналіз конкурентів
    add_content_slide(
        prs,
        "Аналіз конкурентів",
        [
            {'text': 'Sweetwater (США)', 'bold': True, 'size': 18},
            {'text': 'Найбільший онлайн-магазин музичних інструментів', 'level': 1, 'size': 14},
            '',
            {'text': 'Thomann (Німеччина)', 'bold': True, 'size': 18},
            {'text': 'Європейський лідер, понад 1 млн продуктів', 'level': 1, 'size': 14},
            '',
            {'text': 'Guitar Center (США)', 'bold': True, 'size': 18},
            {'text': 'Мережа магазинів + онлайн-платформа', 'level': 1, 'size': 14},
        ]
    )

    # ==================== МЕТОДОЛОГІЯ ====================

    # Slide 6: Методологія Agile/Scrum
    add_content_slide(
        prs,
        "Методологія розробки: Agile/Scrum",
        [
            {'text': 'Agile - гнучка методологія:', 'bold': True, 'size': 18},
            {'text': 'Короткі ітерації (спринти)', 'level': 1},
            {'text': 'Тісна взаємодія з замовником', 'level': 1},
            {'text': 'Швидке реагування на зміни', 'level': 1},
            '',
            {'text': 'Scrum - структура роботи:', 'bold': True, 'size': 18},
            {'text': 'Спринти тривалістю 2-4 тижні', 'level': 1},
            {'text': 'Ефективне планування та коригування', 'level': 1},
            {'text': 'Регулярні retrospective', 'level': 1},
        ]
    )

    # Slide 7: Scrum Sprint Diagram
    if os.path.exists(scrum_diagram):
        add_image_slide(
            prs,
            "Scrum Sprint процес",
            scrum_diagram,
            "Цикл спринту: Product Backlog → Sprint Planning → Sprint → Definition of Done → PSI"
        )

    # Slide 8: Менеджмент проєкту (ClickUp)
    if os.path.exists(clickup_screenshot):
        add_image_slide(
            prs,
            "Менеджмент процесу проєктування",
            clickup_screenshot,
            "Використання ClickUp для управління задачами проєкту (Kanban Board)"
        )

    # ==================== АРХІТЕКТУРА ====================

    # Slide 9: Архітектура
    add_content_slide(
        prs,
        "Архітектура системи",
        [
            {'text': 'Клієнт-серверна архітектура з елементами мікросервісів', 'bold': True, 'size': 18},
            '',
            {'text': 'Клієнтська частина:', 'bold': True},
            {'text': 'Web (Next.js + React)', 'level': 1},
            {'text': 'Mobile (React Native)', 'level': 1},
            '',
            {'text': 'Серверна частина:', 'bold': True},
            {'text': 'RESTful API (NestJS)', 'level': 1},
            {'text': 'PostgreSQL + TypeORM', 'level': 1},
            '',
            {'text': 'Спільні компоненти:', 'bold': True},
            {'text': 'Common Package - типи, валідація, утиліти', 'level': 1},
        ]
    )

    # Slide 10: Technology Stack
    if os.path.exists(tech_stack_img):
        add_image_slide(
            prs,
            "Технологічний стек системи",
            tech_stack_img,
            "Frontend, Backend, Payments & Integrations, DevTools"
        )

    # Slide 11: Package Diagram
    if os.path.exists(package_diagram):
        add_image_slide(
            prs,
            "Діаграма пакетів",
            package_diagram,
            "Monorepo структура проєкту"
        )

    # Slide 12: Class Diagram
    if os.path.exists(class_diagram):
        add_image_slide(
            prs,
            "Діаграма класів",
            class_diagram,
            "Основні сутності системи"
        )

    # ==================== БАЗА ДАНИХ ====================

    # Slide 13: Database Diagram
    if os.path.exists(db_diagram):
        add_image_slide(
            prs,
            "Спроєктована база даних",
            db_diagram,
            "Концептуальна модель бази даних з відношеннями між сутностями"
        )

    # Slide 14: Database Tables
    add_content_slide(
        prs,
        "Структура бази даних - Основні таблиці",
        [
            {'text': 'users - користувачі системи', 'bold': True},
            {'text': 'Поля: id, email, password, firstName, lastName, role, phone', 'level': 1, 'size': 14},
            '',
            {'text': 'products - товари', 'bold': True},
            {'text': 'Поля: id, name, description, price, stock, categoryId', 'level': 1, 'size': 14},
            '',
            {'text': 'orders - замовлення', 'bold': True},
            {'text': 'Поля: id, userId, status, totalAmount, createdAt', 'level': 1, 'size': 14},
            '',
            {'text': 'cart_items, reviews, wishlist, categories', 'bold': True},
            {'text': 'Допоміжні таблиці для функціоналу системи', 'level': 1, 'size': 14},
        ]
    )

    # ==================== ПРОТОТИП ====================

    # Slide 15: Прототип Overview
    if os.path.exists(web_mobile_screens):
        add_image_slide(
            prs,
            "Прототип системи: Web та Mobile",
            web_mobile_screens,
            "Багатоплатформенна реалізація"
        )

    # Slide 16: Веб-застосунок - Каталог
    if os.path.exists(web_catalog):
        add_image_slide(
            prs,
            "Веб-застосунок: Каталог товарів",
            web_catalog,
            "Фільтрація, пошук та сортування товарів"
        )

    # Slide 17: Мобільний застосунок
    if os.path.exists(mobile_catalog):
        add_image_slide(
            prs,
            "Мобільний застосунок: Каталог",
            mobile_catalog,
            "Адаптивний інтерфейс для мобільних пристроїв"
        )

    # Slide 18: Кошик та профіль
    if os.path.exists(cart_screen) and os.path.exists(profile_screen):
        add_two_images_slide(
            prs,
            "Кошик покупок та Профіль користувача",
            cart_screen,
            profile_screen
        )

    # ==================== АДМІНІСТРУВАННЯ ====================

    # Slide 19: User Management
    if os.path.exists(user_management):
        add_image_slide(
            prs,
            "Адмін-панель: Управління користувачами",
            user_management,
            "Перегляд, пошук та управління ролями користувачів"
        )

    # Slide 20: Analytics Dashboard
    if os.path.exists(analytics_dashboard):
        add_image_slide(
            prs,
            "Адмін-панель: Аналітика",
            analytics_dashboard,
            "Dashboard з метриками та статистикою"
        )

    # Slide 21: Функціональність
    add_content_slide(
        prs,
        "Основна функціональність системи",
        [
            {'text': 'Для користувачів:', 'bold': True, 'size': 18},
            {'text': 'Каталог товарів з пошуком та фільтрацією', 'level': 1},
            {'text': 'Детальна інформація про товар', 'level': 1},
            {'text': 'Кошик та оформлення замовлень', 'level': 1},
            {'text': 'Відгуки та рейтинги', 'level': 1},
            {'text': 'Особистий кабінет', 'level': 1},
            '',
            {'text': 'Для адміністраторів:', 'bold': True, 'size': 18},
            {'text': 'Управління товарами та категоріями', 'level': 1},
            {'text': 'Обробка замовлень та користувачів', 'level': 1},
            {'text': 'Аналітика та звіти', 'level': 1},
        ]
    )

    # Slide 22: Безпека
    add_content_slide(
        prs,
        "Безпека системи",
        [
            {'text': 'Автентифікація:', 'bold': True},
            {'text': 'JWT токени з refresh механізмом', 'level': 1, 'size': 14},
            {'text': 'Role-based access control', 'level': 1, 'size': 14},
            '',
            {'text': 'Захист даних:', 'bold': True},
            {'text': 'Bcrypt для паролів', 'level': 1, 'size': 14},
            {'text': 'HTTPS з\'єднання', 'level': 1, 'size': 14},
            {'text': 'Валідація даних (class-validator)', 'level': 1, 'size': 14},
            '',
            {'text': 'Захист від атак:', 'bold': True},
            {'text': 'CORS, Rate limiting, SQL injection prevention', 'level': 1, 'size': 14},
        ]
    )

    # Slide 23: Висновки
    add_content_slide(
        prs,
        "Висновки",
        [
            {'text': 'Досягнуті результати:', 'bold': True, 'size': 18},
            '',
            "Проаналізовано методи дослідження та принципи UI/UX",
            "Досліджено ринок та сформовано вимоги",
            "Обрано методологію Agile/Scrum",
            "Використано ClickUp для менеджменту проєкту",
            "Спроектовано архітектуру та базу даних",
            "Обрано технологічний стек",
            "Реалізовано веб та мобільні застосунки",
            "Впроваджено адмін-панель з аналітикою",
            "Забезпечено безпеку системи"
        ]
    )

    # Save presentation
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    try:
        output = create_presentation()
        print("SUCCESS")
    except Exception as e:
        print("ERROR")
        import traceback
        traceback.print_exc()
